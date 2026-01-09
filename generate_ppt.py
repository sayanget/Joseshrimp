from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

# Configuración de archivos
SCREENSHOTS_DIR = "d:/project/jose/screenshots"
OUTPUT_FILE = "d:/project/jose/MANUAL_USUARIO_ES.pptx"

# Mapeo de diapositivas a imágenes (ajustar si los nombres cambian levemente)
def get_screenshot(prefix):
    files = [f for f in os.listdir(SCREENSHOTS_DIR) if f.startswith(prefix)]
    if files:
        return os.path.join(SCREENSHOTS_DIR, files[0])
    return None

def create_presentation():
    prs = Presentation()

    # --- 1. Portada ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Sistema de Gestión de Ventas"
    subtitle.text = "Manual de Usuario y Guía de Operación\nJose Burgueno - Enero 2026"

    # --- 2. Introducción ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Introducción al Sistema"
    content = slide.placeholders[1]
    content.text = (
        "- Aplicación web moderna para gestión de ventas e inventario.\n"
        "- Enfoque en trazabilidad, precisión y facilidad de uso.\n"
        "- Multilingüe: Soporte completo para Español, Inglés y Chino.\n"
        "- Principio clave: Cálculo automático de peso y prohibición de entrada libre."
    )

    # --- 3. Panel Principal (Dashboard) ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Panel Principal (Dashboard)"
    
    img_path = get_screenshot("screenshot_dashboard_es")
    if img_path:
        prs.slides[2].shapes.add_picture(img_path, Inches(0.5), Inches(1.5), height=Inches(4.5))
    
    left = Inches(5.5)
    top = Inches(1.5)
    width = Inches(4)
    height = Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Vista General del Negocio"
    p = tf.add_paragraph()
    p.text = "- Resumen de pedidos y ventas del día.\n- Estado del stock y alertas críticas.\n- Acceso rápido a funciones principales."
    p.font.size = Pt(18)

    # --- 4. Gestión de Ventas (Lista) ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Gestión de Ventas: Lista"
    
    img_path = get_screenshot("screenshot_sales_es")
    if img_path:
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), height=Inches(4.5))
        
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Control de Órdenes"
    p = tf.add_paragraph()
    p.text = "- Listado histórico de todas las ventas.\n- Filtrado por cliente, fecha y estado.\n- Acceso a detalles y anulación de órdenes."
    p.font.size = Pt(18)

    # --- 5. Nueva Venta ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Creación de Nueva Venta"
    
    img_path = get_screenshot("screenshot_new_sale_es")
    if img_path:
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), height=Inches(4.5))
        
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Operación Ágil y Segura"
    p = tf.add_paragraph()
    p.text = "- Selección de cliente y tipo de pago.\n- Agregado dinámico de productos.\n- Cálculo automático de pesos (Cajas x KG + Extra).\n- Validación de crédito en tiempo real."
    p.font.size = Pt(18)

    # --- 6. Inventario (Stock Actual) ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Control de Inventario"
    
    img_path = get_screenshot("screenshot_inventory_es")
    if img_path:
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), height=Inches(4.5))
        
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Visibilidad Total del Stock"
    p = tf.add_paragraph()
    p.text = "- Monitoreo de existencias en KG.\n- Código de colores para alertas de stock.\n- Historial completo de movimientos (entradas/salidas)."
    p.font.size = Pt(18)

    # --- 7. Reportes y Estadísticas ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Reportes y Estadísticas"
    
    img_path = get_screenshot("screenshot_report_es")
    if img_path:
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), height=Inches(4.5))
        
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Toma de Decisiones"
    p = tf.add_paragraph()
    p.text = "- Gráficos interactivos de ventas diarias.\n- Análisis por cliente y por producto.\n- Análisis de proporción de producto suelto."
    p.font.size = Pt(18)

    # --- 8. Administración (Especificaciones) ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Panel de Administración"
    
    img_path = get_screenshot("screenshot_admin_es")
    if img_path:
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), height=Inches(4.5))
        
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Configuración del Sistema"
    p = tf.add_paragraph()
    p.text = "- Gestión de especificaciones de productos.\n- Administración de clientes y permisos de crédito.\n- Auditoría completa de todas las acciones."
    p.font.size = Pt(18)

    # --- 9. Conclusión ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Gracias"
    content = slide.placeholders[1]
    content.text = (
        "El sistema está diseñado para crecer con su negocio.\n\n"
        "Para soporte técnico o dudas, contacte al administrador.\n\n"
        "© 2026 Sistema de Gestión de Ventas"
    )

    # Guardar
    prs.save(OUTPUT_FILE)
    print(f"✅ Presentación generada exitosamente: {OUTPUT_FILE}")

if __name__ == "__main__":
    create_presentation()
